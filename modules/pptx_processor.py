"""
PowerPoint Processing Module

Learning Note - What This Module Does:
--------------------------------------
This module extracts content from PowerPoint (.pptx) files for the RAG system.

PowerPoint files contain:
1. Slide text (titles, body, bullet points)
2. Speaker notes (presenter comments)
3. Embedded images (charts, diagrams, photos)
4. Tables and shapes

Each slide becomes a searchable chunk with full context preserved.

Key Features:
- Extract text from all slides
- Extract speaker notes
- Extract and analyze embedded images (using existing vision processor)
- Maintain slide order and structure
- Rich metadata (slide number, title, layout)

Why PowerPoint Support Matters:
- Presentations contain condensed, important information
- Visual data (charts, diagrams) is critical to understanding
- Speaker notes provide additional context not visible in slides
- Many educational and business materials are in PowerPoint format
"""

from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from io import BytesIO
import base64
from config import settings


class PowerPointProcessor:
    """
    Processes PowerPoint files to extract text, images, and metadata.

    Learning Note - Processing Strategy:
    ------------------------------------
    1. Open presentation and iterate through slides
    2. Extract title and body text from each slide
    3. Extract speaker notes if enabled
    4. Extract embedded images if enabled
    5. Create structured chunks with metadata
    6. Preserve slide context (don't split slides across chunks)
    """

    def __init__(self):
        """Initialize the PowerPoint processor."""
        self.max_slides = settings.MAX_SLIDES_PER_DOCUMENT
        self.include_notes = settings.INCLUDE_SPEAKER_NOTES
        self.extract_images = settings.EXTRACT_SLIDE_IMAGES and settings.ENABLE_IMAGE_PROCESSING
        self.preserve_structure = settings.PRESERVE_SLIDE_STRUCTURE

    def process_presentation(self, pptx_path: str) -> Dict[str, Any]:
        """
        Process a PowerPoint file and extract all content.

        Args:
            pptx_path: Path to the .pptx file

        Returns:
            Dictionary containing:
            - slides: List of slide data (text, images, metadata)
            - total_slides: Total number of slides
            - processed_slides: Number of slides actually processed
            - images_extracted: Number of images extracted

        Learning Note - Why Return Structured Data?
        -------------------------------------------
        We return structured data rather than flat chunks because:
        - Allows flexible chunk creation in document_processor
        - Preserves slide relationships
        - Enables different chunking strategies
        - Makes debugging and testing easier
        """
        try:
            print(f"Opening PowerPoint: {pptx_path}")
            presentation = Presentation(pptx_path)

            total_slides = len(presentation.slides)
            slides_to_process = min(total_slides, self.max_slides)

            print(f"Found {total_slides} slides, processing {slides_to_process}")

            slides_data = []
            images_extracted = 0

            for slide_num, slide in enumerate(presentation.slides[:slides_to_process], 1):
                slide_data = self._process_slide(slide, slide_num, total_slides)
                slides_data.append(slide_data)

                if slide_data.get('images'):
                    images_extracted += len(slide_data['images'])

                if slide_num % 10 == 0:
                    print(f"Processed {slide_num}/{slides_to_process} slides...")

            print(f"✓ Extraction complete: {slides_to_process} slides, {images_extracted} images")

            return {
                'slides': slides_data,
                'total_slides': total_slides,
                'processed_slides': slides_to_process,
                'images_extracted': images_extracted
            }

        except Exception as e:
            print(f"Error processing PowerPoint: {e}")
            return {
                'slides': [],
                'total_slides': 0,
                'processed_slides': 0,
                'images_extracted': 0,
                'error': str(e)
            }

    def _process_slide(self, slide, slide_num: int, total_slides: int) -> Dict[str, Any]:
        """
        Process a single slide to extract all content.

        Args:
            slide: python-pptx Slide object
            slide_num: Slide number (1-indexed)
            total_slides: Total slides in presentation

        Returns:
            Dictionary with slide data

        Learning Note - Slide Structure:
        --------------------------------
        PowerPoint slides have:
        - Title shape (optional)
        - Body shapes (text boxes, bullet points)
        - Notes (speaker notes)
        - Images (embedded pictures)
        - Tables, charts, etc.

        We extract text from all shapes and combine into structured data.
        """
        slide_data = {
            'slide_number': slide_num,
            'total_slides': total_slides,
            'title': '',
            'body': '',
            'notes': '',
            'images': [],
            'layout': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown'
        }

        # Extract title
        if slide.shapes.title:
            slide_data['title'] = slide.shapes.title.text.strip()

        # Extract body text from all text shapes
        body_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape != slide.shapes.title:
                text = shape.text.strip()
                if text:
                    body_parts.append(text)

        slide_data['body'] = '\n'.join(body_parts)

        # Extract speaker notes if enabled
        if self.include_notes and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if hasattr(notes_slide, 'notes_text_frame'):
                slide_data['notes'] = notes_slide.notes_text_frame.text.strip()

        # Extract images if enabled
        if self.extract_images:
            slide_data['images'] = self._extract_slide_images(slide, slide_num)

        return slide_data

    def _extract_slide_images(self, slide, slide_num: int) -> List[Dict[str, Any]]:
        """
        Extract embedded images from a slide.

        Args:
            slide: python-pptx Slide object
            slide_num: Slide number for metadata

        Returns:
            List of image data dictionaries

        Learning Note - Image Extraction:
        ---------------------------------
        Images in slides can be:
        - Pictures (photos, diagrams)
        - Charts (converted to images)
        - Screenshots
        - Logos and icons

        We extract the image binary data and metadata for later vision analysis.
        """
        images = []
        image_count = 0

        for shape in slide.shapes:
            try:
                # Check if shape is a picture
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1

                    # Get image binary data
                    image = shape.image
                    image_bytes = image.blob

                    # Load with PIL to get dimensions
                    pil_image = Image.open(BytesIO(image_bytes))
                    width, height = pil_image.size

                    # Filter by size (skip small icons/decorations)
                    if width >= settings.MIN_IMAGE_WIDTH and height >= settings.MIN_IMAGE_HEIGHT:
                        # Convert to base64 for storage
                        image_base64 = base64.b64encode(image_bytes).decode('utf-8')

                        images.append({
                            'image_data': image_base64,
                            'width': width,
                            'height': height,
                            'format': image.content_type.split('/')[-1],
                            'slide_number': slide_num,
                            'image_index': image_count
                        })
                    else:
                        print(f"  Skipping small image on slide {slide_num}: {width}x{height}")

            except Exception as e:
                print(f"Error extracting image from slide {slide_num}: {e}")
                continue

        return images

    def create_slide_chunks(self, slides_data: List[Dict[str, Any]], filename: str) -> List[Dict[str, Any]]:
        """
        Create searchable chunks from slide data.

        Args:
            slides_data: List of slide dictionaries from process_presentation()
            filename: Original filename for metadata

        Returns:
            List of chunks ready for vector storage

        Learning Note - Chunking Strategy:
        ----------------------------------
        For PowerPoint, we have two strategies:

        1. PRESERVE_SLIDE_STRUCTURE=True (recommended):
           - Each slide = one chunk (title + body + notes together)
           - Preserves context and makes results more useful
           - Easier to cite ("See slide 5")

        2. PRESERVE_SLIDE_STRUCTURE=False:
           - Content split by PPTX_CHUNK_SIZE
           - May split slides across chunks
           - Only useful for very large slides

        Most presentations work best with strategy #1.
        """
        chunks = []

        if self.preserve_structure:
            # Strategy 1: Each slide is a chunk
            for slide in slides_data:
                chunk_text = self._format_slide_text(slide)

                if chunk_text.strip():
                    chunk = {
                        'text': chunk_text,
                        'type': 'slide',
                        'metadata': {
                            'source': filename,
                            'chunk_type': 'slide',
                            'slide_number': slide['slide_number'],
                            'total_slides': slide['total_slides'],
                            'slide_title': slide['title'],
                            'layout': slide['layout'],
                            'has_images': len(slide['images']) > 0,
                            'num_images': len(slide['images'])
                        }
                    }
                    chunks.append(chunk)

        else:
            # Strategy 2: Split by chunk size (less common)
            all_text = []
            for slide in slides_data:
                slide_text = self._format_slide_text(slide)
                all_text.append(slide_text)

            combined_text = '\n\n'.join(all_text)

            # Split into chunks of PPTX_CHUNK_SIZE
            chunk_size = settings.PPTX_CHUNK_SIZE
            for i in range(0, len(combined_text), chunk_size):
                chunk_text = combined_text[i:i + chunk_size]

                if chunk_text.strip():
                    chunk = {
                        'text': chunk_text,
                        'type': 'text',
                        'metadata': {
                            'source': filename,
                            'chunk_type': 'pptx_text',
                            'chunk_index': i // chunk_size
                        }
                    }
                    chunks.append(chunk)

        return chunks

    def _format_slide_text(self, slide: Dict[str, Any]) -> str:
        """
        Format slide content into readable text.

        Args:
            slide: Slide data dictionary

        Returns:
            Formatted text string

        Learning Note - Text Formatting:
        --------------------------------
        We format slide content to be human-readable:
        - Clear section headers
        - Proper spacing
        - Slide number for context
        - All text components included

        This makes search results more useful to users.
        """
        parts = []

        # Add slide header
        if slide['title']:
            parts.append(f"=== Slide {slide['slide_number']}: {slide['title']} ===")
        else:
            parts.append(f"=== Slide {slide['slide_number']} ===")

        # Add body content
        if slide['body']:
            parts.append(slide['body'])

        # Add speaker notes if present
        if slide['notes']:
            parts.append(f"\n[Speaker Notes]")
            parts.append(slide['notes'])

        # Add image indicator
        if slide['images']:
            parts.append(f"\n[Contains {len(slide['images'])} image(s)]")

        return '\n\n'.join(parts)
