"""
Document Processing Module for RAG Agent.

This module handles:
1. Parsing different document formats (PDF, DOCX, TXT, PPTX)
2. Intelligent text chunking with overlap
3. Metadata attachment for source tracking
4. Image extraction and analysis

Learning Note:
Document processing is the first step in the RAG pipeline. Good chunking
is crucial because:
- It determines how precisely we can retrieve information
- Poor chunking can split important context
- Proper overlap maintains continuity between chunks
"""

import re
from datetime import datetime
from typing import List, Dict, Any, BinaryIO
from pathlib import Path
import tempfile
import shutil

# Document parsing libraries
import pypdf
from docx import Document
import chardet

# Configuration
from config import settings

# Image processing (conditional import)
try:
    from modules.image_extractor import ImageExtractor
    from modules.vision_processor import VisionProcessor
    IMAGE_PROCESSING_AVAILABLE = True
except ImportError:
    IMAGE_PROCESSING_AVAILABLE = False
    print("Image processing modules not available.")

# PowerPoint processing (conditional import)
try:
    from modules.pptx_processor import PowerPointProcessor
    PPTX_PROCESSING_AVAILABLE = True
except ImportError:
    PPTX_PROCESSING_AVAILABLE = False
    print("PowerPoint processing module not available.")


class DocumentParser:
    """
    Parses documents of various formats and extracts text.

    Supported formats:
    - PDF (.pdf) - using pypdf
    - Microsoft Word (.docx) - using python-docx
    - Microsoft PowerPoint (.pptx) - using python-pptx
    - Plain text (.txt) - with automatic encoding detection
    """

    @staticmethod
    def parse_pdf(file: BinaryIO) -> str:
        """
        Extract text from a PDF file.

        Args:
            file: Binary file object (PDF)

        Returns:
            str: Extracted text from all pages

        Learning Note:
            PDF parsing can be tricky because:
            - PDFs can contain images, tables, and complex layouts
            - Text extraction isn't always perfect
            - Some PDFs are scanned images (would need OCR)
            - Password-protected PDFs need special handling
        """
        try:
            # Create PDF reader object
            pdf_reader = pypdf.PdfReader(file)

            # Extract text from all pages
            text_parts = []
            for page_num, page in enumerate(pdf_reader.pages, start=1):
                page_text = page.extract_text()

                if page_text.strip():  # Only add non-empty pages
                    text_parts.append(page_text)

            # Combine all pages
            full_text = "\n\n".join(text_parts)

            if not full_text.strip():
                raise ValueError("PDF appears to be empty or contains only images")

            return full_text

        except pypdf.errors.PdfReadError as e:
            raise ValueError(f"Could not read PDF file: {str(e)}")
        except Exception as e:
            raise ValueError(f"Error parsing PDF: {str(e)}")

    @staticmethod
    def parse_docx(file: BinaryIO) -> str:
        """
        Extract text from a Microsoft Word document.

        Args:
            file: Binary file object (DOCX)

        Returns:
            str: Extracted text from all paragraphs

        Learning Note:
            DOCX files are actually ZIP archives containing XML files.
            The python-docx library handles this complexity for us.
            It extracts text from paragraphs but may miss:
            - Headers and footers
            - Text in text boxes
            - Text in tables (though we extract these)
        """
        try:
            # Load document
            doc = Document(file)

            # Extract text from all paragraphs
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

            # Also extract text from tables
            table_texts = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        table_texts.append(row_text)

            # Combine paragraphs and tables
            all_text = paragraphs + table_texts
            full_text = "\n\n".join(all_text)

            if not full_text.strip():
                raise ValueError("Document appears to be empty")

            return full_text

        except Exception as e:
            raise ValueError(f"Error parsing DOCX: {str(e)}")

    @staticmethod
    def parse_txt(file: BinaryIO) -> str:
        """
        Read a plain text file with automatic encoding detection.

        Args:
            file: Binary file object (TXT)

        Returns:
            str: Text content

        Learning Note:
            Text files can have different encodings:
            - UTF-8 (most common, supports all languages)
            - ASCII (English only, subset of UTF-8)
            - Latin-1, Windows-1252, etc.

            The chardet library detects encoding automatically,
            preventing errors with non-UTF-8 files.
        """
        try:
            # Read raw bytes
            raw_bytes = file.read()

            # Detect encoding
            detected = chardet.detect(raw_bytes)
            encoding = detected.get('encoding', 'utf-8')

            # Fallback to utf-8 if detection failed
            if not encoding:
                encoding = 'utf-8'

            # Decode with detected encoding
            text = raw_bytes.decode(encoding, errors='replace')

            if not text.strip():
                raise ValueError("Text file appears to be empty")

            return text

        except Exception as e:
            raise ValueError(f"Error parsing text file: {str(e)}")

    @staticmethod
    def parse_pptx(file: BinaryIO) -> str:
        """
        Extract text from a PowerPoint file.

        Args:
            file: Binary file object (PPTX)

        Returns:
            str: Extracted text from all slides

        Learning Note - PowerPoint vs Other Formats:
        --------------------------------------------
        PowerPoint files have unique structure:
        - Slides with titles and body text
        - Speaker notes (often ignored but valuable!)
        - Embedded images (need separate extraction)
        - Tables and charts

        We extract text-only here. Full slide processing with images
        is handled in process_document() method using PowerPointProcessor.

        This method is used when ENABLE_PPTX_PROCESSING=False or
        when only text extraction is needed (no slide structure preservation).
        """
        try:
            from pptx import Presentation

            # Load presentation
            presentation = Presentation(file)

            # Extract text from all slides
            text_parts = []

            for slide_num, slide in enumerate(presentation.slides, start=1):
                slide_parts = []

                # Extract title
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                    if title:
                        slide_parts.append(f"Slide {slide_num}: {title}")

                # Extract body text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape != slide.shapes.title:
                        text = shape.text.strip()
                        if text:
                            slide_parts.append(text)

                # Extract speaker notes if available
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if hasattr(notes_slide, 'notes_text_frame'):
                        notes_text = notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_parts.append(f"Notes: {notes_text}")

                # Combine slide parts
                if slide_parts:
                    text_parts.append('\n'.join(slide_parts))

            # Combine all slides
            full_text = "\n\n".join(text_parts)

            if not full_text.strip():
                raise ValueError("PowerPoint appears to be empty or contains only images")

            return full_text

        except Exception as e:
            raise ValueError(f"Error parsing PowerPoint file: {str(e)}")

    @staticmethod
    def parse_document(file: BinaryIO, file_type: str) -> str:
        """
        Route document to appropriate parser based on file type.

        Args:
            file: Binary file object
            file_type: File extension (e.g., '.pdf', '.docx', '.txt')

        Returns:
            str: Extracted text

        Raises:
            ValueError: If file type is not supported or parsing fails
        """
        file_type = file_type.lower()

        if file_type == '.pdf':
            return DocumentParser.parse_pdf(file)
        elif file_type == '.docx':
            return DocumentParser.parse_docx(file)
        elif file_type == '.pptx':
            return DocumentParser.parse_pptx(file)
        elif file_type == '.txt':
            return DocumentParser.parse_txt(file)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")


class TextChunker:
    """
    Intelligently chunks text into smaller pieces with overlap.

    Why chunking matters:
    - LLMs have limited context windows
    - Smaller chunks = more precise retrieval
    - Overlap ensures context isn't lost at boundaries
    """

    @staticmethod
    def chunk_by_sentences(
        text: str,
        chunk_size: int = None,
        overlap: int = None
    ) -> List[str]:
        """
        Split text into chunks, trying to preserve sentence boundaries.

        Args:
            text: Text to chunk
            chunk_size: Target size of each chunk in characters
            overlap: Number of characters to overlap between chunks

        Returns:
            List of text chunks

        Learning Note:
            Sentence-aware chunking is better than simple character splitting because:
            - Complete sentences make more semantic sense
            - Easier for embeddings to capture meaning
            - Better for human readability
            - Prevents awkward mid-sentence cuts

        Algorithm:
        1. Split text into sentences
        2. Accumulate sentences until chunk_size is reached
        3. Start new chunk with overlap from previous chunk
        4. Repeat until all text is processed
        """
        if chunk_size is None:
            chunk_size = settings.CHUNK_SIZE

        if overlap is None:
            overlap = settings.CHUNK_OVERLAP

        # Split text into sentences using regex
        # This pattern looks for sentence-ending punctuation followed by space
        sentence_pattern = r'(?<=[.!?])\s+'
        sentences = re.split(sentence_pattern, text)

        chunks = []
        current_chunk = ""

        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue

            # Check if adding this sentence would exceed chunk_size
            potential_chunk = current_chunk + " " + sentence if current_chunk else sentence

            if len(potential_chunk) > chunk_size and current_chunk:
                # Current chunk is full, save it and start new one
                chunks.append(current_chunk.strip())

                # Start new chunk with overlap from previous chunk
                # Take the last 'overlap' characters from current chunk
                overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk

                # Find the start of the first complete word in overlap
                # (avoid starting mid-word)
                space_idx = overlap_text.find(' ')
                if space_idx != -1:
                    overlap_text = overlap_text[space_idx + 1:]

                current_chunk = overlap_text + " " + sentence
            else:
                # Add sentence to current chunk
                current_chunk = potential_chunk

        # Don't forget the last chunk
        if current_chunk.strip():
            chunks.append(current_chunk.strip())

        return chunks

    @staticmethod
    def add_metadata(
        chunks: List[str],
        filename: str,
        file_type: str
    ) -> List[Dict[str, Any]]:
        """
        Attach metadata to each chunk for source tracking.

        Args:
            chunks: List of text chunks
            filename: Source filename
            file_type: File extension (e.g., '.pdf')

        Returns:
            List of dictionaries with 'text' and 'metadata' keys

        Learning Note:
            Metadata is crucial for:
            - Source attribution (which document/chunk did this come from?)
            - Debugging (why was this chunk retrieved?)
            - Filtering (search within specific documents)
            - Timestamps (when was this uploaded?)

        Metadata schema:
        {
            'filename': str - original filename
            'chunk_id': int - position in document (0-indexed)
            'file_type': str - file extension
            'upload_date': str - ISO format timestamp
            'total_chunks': int - how many chunks in this document
        }
        """
        chunks_with_metadata = []

        for idx, chunk_text in enumerate(chunks):
            chunk_dict = {
                'text': chunk_text,
                'metadata': {
                    'filename': filename,
                    'chunk_id': idx,
                    'file_type': file_type,
                    'upload_date': datetime.now().isoformat(),
                    'total_chunks': len(chunks)
                }
            }
            chunks_with_metadata.append(chunk_dict)

        return chunks_with_metadata

    @staticmethod
    def process_document(
        file: BinaryIO,
        filename: str,
        file_type: str
    ) -> List[Dict[str, Any]]:
        """
        Complete pipeline: Parse document, chunk text, add metadata, process images.

        Args:
            file: Binary file object
            filename: Original filename
            file_type: File extension

        Returns:
            List of chunks with metadata (includes both text and image chunks)

        This is the main entry point for document processing.
        It orchestrates the entire pipeline:
        1. Parse document to extract text
        2. Chunk text intelligently
        3. Extract and analyze images (if enabled)
        4. Add metadata for tracking
        5. Combine text and image chunks

        Learning Note - Image Processing Flow:
        -------------------------------------
        When image processing is enabled:
        1. Save uploaded file temporarily (image extraction needs file path)
        2. Extract images from document
        3. Analyze each image with GPT-4 Vision
        4. Create special "image chunks" with descriptions
        5. Clean up temporary file
        6. Combine text chunks and image chunks

        Image chunks have:
        - type: 'image'
        - text: Image description from vision model
        - image_data: Base64 encoded image
        - metadata: Image info (size, page, etc.)

        Example:
            >>> with open('document.pdf', 'rb') as f:
            >>>     chunks = TextChunker.process_document(f, 'document.pdf', '.pdf')
            >>> print(f"Created {len(chunks)} chunks")
            >>> text_chunks = [c for c in chunks if c.get('type') != 'image']
            >>> image_chunks = [c for c in chunks if c.get('type') == 'image']
            >>> print(f"Text: {len(text_chunks)}, Images: {len(image_chunks)}")
        """
        # Special handling for PowerPoint when full processing is enabled
        if file_type.lower() == '.pptx' and settings.ENABLE_PPTX_PROCESSING and PPTX_PROCESSING_AVAILABLE:
            print(f"\nProcessing PowerPoint with full slide structure preservation...")
            return TextChunker._process_powerpoint_full(file, filename)

        # Step 1: Parse document to get text
        file.seek(0)  # Reset file pointer
        text = DocumentParser.parse_document(file, file_type)

        # Step 2: Chunk the text
        chunks = TextChunker.chunk_by_sentences(text)

        # Step 3: Add metadata
        chunks_with_metadata = TextChunker.add_metadata(
            chunks,
            filename,
            file_type
        )

        # Step 4: Process images (if enabled and supported)
        image_chunks = []
        if settings.ENABLE_IMAGE_PROCESSING and IMAGE_PROCESSING_AVAILABLE:
            # Check if file type supports image extraction
            if file_type.lower() in ['.pdf', '.docx', '.doc', '.pptx']:
                print(f"\nProcessing images from {filename}...")

                # Save file temporarily for image extraction
                temp_file = None
                try:
                    # Create temporary file
                    file.seek(0)  # Reset file pointer
                    with tempfile.NamedTemporaryFile(
                        delete=False,
                        suffix=file_type,
                        dir=settings.IMAGE_TEMP_DIR
                    ) as temp_file:
                        # Copy uploaded file to temporary file
                        shutil.copyfileobj(file, temp_file)
                        temp_path = temp_file.name

                    # Extract images
                    extractor = ImageExtractor()
                    images = extractor.extract_images(temp_path, file_type.lstrip('.'))

                    if images:
                        print(f"Found {len(images)} images")

                        # Analyze images with vision model
                        processor = VisionProcessor()
                        analyzed_images = processor.process_images_batch(images, show_progress=True)

                        # Create chunks for images
                        for img in analyzed_images:
                            image_chunk = {
                                'text': f"[IMAGE] {img['description']}",
                                'type': 'image',
                                'image_data': img['image_data'],
                                'image_format': img['format'],
                                'metadata': {
                                    'filename': filename,
                                    'file_type': file_type,
                                    'upload_date': datetime.now().isoformat(),
                                    'chunk_type': 'image',
                                    'image_width': img['width'],
                                    'image_height': img['height'],
                                    'image_source': img['source_type'],
                                    'vision_model': img.get('vision_model', settings.VISION_MODEL),
                                    'page_number': img.get('page_number', 'unknown')
                                }
                            }
                            image_chunks.append(image_chunk)

                        print(f"Created {len(image_chunks)} image chunks")
                    else:
                        print("No images found in document")

                except Exception as e:
                    print(f"Error processing images: {e}")
                finally:
                    # Clean up temporary file
                    if temp_file and Path(temp_path).exists():
                        try:
                            Path(temp_path).unlink()
                        except Exception as e:
                            print(f"Could not delete temp file: {e}")

        # Step 5: Combine text and image chunks
        all_chunks = chunks_with_metadata + image_chunks

        print(f"\nTotal chunks created: {len(all_chunks)} ({len(chunks_with_metadata)} text, {len(image_chunks)} images)")

        return all_chunks

    @staticmethod
    def _process_powerpoint_full(file: BinaryIO, filename: str) -> List[Dict[str, Any]]:
        """
        Process PowerPoint with full slide structure preservation.

        Args:
            file: Binary file object (PPTX)
            filename: Original filename

        Returns:
            List of chunks with slide structure and images

        Learning Note - Full PowerPoint Processing:
        -------------------------------------------
        This method provides complete PowerPoint processing:
        1. Extract text from each slide (title, body, notes)
        2. Preserve slide structure (each slide = one chunk)
        3. Extract and analyze images from slides
        4. Create rich metadata (slide numbers, titles, etc.)

        This is different from generic text extraction because:
        - Slides are natural chunk boundaries
        - Slide numbers and titles provide context
        - Speaker notes add valuable information
        - Slide images are part of the content
        """
        try:
            # Save file temporarily for processing
            temp_file = None
            try:
                # Create temporary file
                file.seek(0)
                with tempfile.NamedTemporaryFile(
                    delete=False,
                    suffix='.pptx',
                    dir=settings.IMAGE_TEMP_DIR
                ) as temp_file:
                    shutil.copyfileobj(file, temp_file)
                    temp_path = temp_file.name

                # Process PowerPoint
                pptx_processor = PowerPointProcessor()
                presentation_data = pptx_processor.process_presentation(temp_path)

                if 'error' in presentation_data:
                    raise ValueError(presentation_data['error'])

                # Create slide chunks
                slide_chunks = pptx_processor.create_slide_chunks(
                    presentation_data['slides'],
                    filename
                )

                # Process slide images if enabled
                image_chunks = []
                if settings.EXTRACT_SLIDE_IMAGES and IMAGE_PROCESSING_AVAILABLE:
                    print(f"\nAnalyzing {presentation_data['images_extracted']} slide images...")

                    # Collect all images from all slides
                    all_images = []
                    for slide in presentation_data['slides']:
                        all_images.extend(slide['images'])

                    if all_images:
                        # Analyze with vision processor
                        vision_processor = VisionProcessor()
                        analyzed_images = vision_processor.process_images_batch(all_images)

                        # Create image chunks
                        for img in analyzed_images:
                            image_chunk = {
                                'text': f"[IMAGE from Slide {img['slide_number']}] {img['description']}",
                                'type': 'image',
                                'image_data': img['image_data'],
                                'metadata': {
                                    'source': filename,
                                    'chunk_type': 'image',
                                    'source_type': 'pptx_slide',
                                    'slide_number': img['slide_number'],
                                    'image_width': img['width'],
                                    'image_height': img['height'],
                                    'image_index': img.get('image_index', 0),
                                    'vision_model': img.get('vision_model', settings.VISION_MODEL)
                                }
                            }
                            image_chunks.append(image_chunk)

                        print(f"Created {len(image_chunks)} image chunks from slides")

                # Combine slide and image chunks
                all_chunks = slide_chunks + image_chunks

                print(f"\nTotal PowerPoint chunks: {len(all_chunks)} ({len(slide_chunks)} slides, {len(image_chunks)} images)")

                return all_chunks

            finally:
                # Clean up temporary file
                if temp_file and Path(temp_path).exists():
                    try:
                        Path(temp_path).unlink()
                    except Exception as e:
                        print(f"Could not delete temp file: {e}")

        except Exception as e:
            print(f"Error in full PowerPoint processing: {e}")
            # Fallback to basic text extraction
            print("Falling back to basic text extraction...")
            file.seek(0)
            text = DocumentParser.parse_pptx(file)
            chunks = TextChunker.chunk_by_sentences(text)
            return TextChunker.add_metadata(chunks, filename, '.pptx')
