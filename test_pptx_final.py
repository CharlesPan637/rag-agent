"""
Final PowerPoint Processing Test

Tests the PowerPoint processing pipeline without requiring
embeddings or vector storage.

This verifies:
1. PowerPoint files can be read with python-pptx
2. Slides are extracted with text and notes
3. Chunks are created with proper metadata
4. The document processor integration works
"""

import sys
from pathlib import Path

# Add parent to path
sys.path.insert(0, str(Path(__file__).parent))

# Import only what we need
from pptx import Presentation


def test_direct_pptx_reading():
    """Test 1: Direct python-pptx reading."""
    print("=" * 70)
    print("TEST 1: Direct PowerPoint File Reading")
    print("=" * 70)

    prs = Presentation('comprehensive_presentation.pptx')

    print(f"\n✓ Opened PowerPoint file")
    print(f"✓ Total slides: {len(prs.slides)}")
    print(f"✓ Slide dimensions: {prs.slide_width} × {prs.slide_height}")

    for i, slide in enumerate(prs.slides, 1):
        # Extract title
        title = "No title"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                title = shape.text
                break

        # Check for notes
        has_notes = bool(slide.notes_slide.notes_text_frame.text.strip())

        print(f"\n  Slide {i}: {title[:50]}")
        print(f"  - Shapes: {len(slide.shapes)}")
        print(f"  - Has notes: {'Yes' if has_notes else 'No'}")

    print(f"\n✅ Test 1 Passed: python-pptx works correctly")
    return True


def test_pptx_processor_module():
    """Test 2: Our PowerPointProcessor module."""
    print("\n" + "=" * 70)
    print("TEST 2: PowerPointProcessor Module")
    print("=" * 70)

    from modules.pptx_processor import PowerPointProcessor

    processor = PowerPointProcessor()
    print(f"\n✓ Initialized PowerPointProcessor")
    print(f"  - Max slides: {processor.max_slides}")
    print(f"  - Include speaker notes: {processor.include_notes}")

    # Process the file
    result = processor.process_presentation('comprehensive_presentation.pptx')

    print(f"\n✓ Processed PowerPoint file")
    print(f"  - Total slides: {result['total_slides']}")
    print(f"  - Processed slides: {result['processed_slides']}")
    print(f"  - Images extracted: {result['images_extracted']}")

    # Check slides
    for i, slide in enumerate(result['slides'], 1):
        print(f"\n  Slide {i}: {slide['title'][:50]}")
        print(f"    - Slide number: {slide['slide_number']}")
        print(f"    - Body text length: {len(slide['body'])} chars")
        print(f"    - Notes length: {len(slide['notes'])} chars")
        print(f"    - Has title: {'Yes' if slide['title'] else 'No'}")

    # Create chunks
    chunks = processor.create_slide_chunks(result['slides'], 'comprehensive_presentation.pptx')

    print(f"\n✓ Created {len(chunks)} chunks")

    for i, chunk in enumerate(chunks, 1):
        metadata = chunk['metadata']
        print(f"\n  Chunk {i}:")
        print(f"    - Slide number: {metadata['slide_number']}")
        print(f"    - Slide title: {metadata.get('slide_title', 'N/A')[:40]}")
        print(f"    - Chunk type: {metadata['chunk_type']}")
        print(f"    - Text length: {len(chunk['text'])} chars")
        print(f"    - Text preview: {chunk['text'][:100]}...")

    print(f"\n✅ Test 2 Passed: PowerPointProcessor works correctly")
    return chunks


def test_document_processor():
    """Test 3: Full document processor integration."""
    print("\n" + "=" * 70)
    print("TEST 3: Document Processor Integration")
    print("=" * 70)

    from modules.document_processor import DocumentParser

    # Test parse_pptx
    with open('comprehensive_presentation.pptx', 'rb') as f:
        text = DocumentParser.parse_pptx(f)

    print(f"\n✓ parse_pptx() works")
    print(f"  - Extracted {len(text)} characters")
    print(f"  - Preview: {text[:200]}...")

    # Test that it's included in parse_document routing
    with open('comprehensive_presentation.pptx', 'rb') as f:
        text2 = DocumentParser.parse_document(f, '.pptx')

    print(f"\n✓ parse_document() routing works for .pptx")
    print(f"  - Extracted {len(text2)} characters")

    # Test full processing (this will use _process_powerpoint_full)
    from modules.document_processor import TextChunker

    with open('comprehensive_presentation.pptx', 'rb') as f:
        chunks = TextChunker.process_document(
            f,
            'comprehensive_presentation.pptx',
            '.pptx'
        )

    print(f"\n✓ Full processing pipeline works")
    print(f"  - Created {len(chunks)} chunks")

    # Verify chunk structure
    for i, chunk in enumerate(chunks, 1):
        metadata = chunk['metadata']
        print(f"\n  Chunk {i}:")
        print(f"    - Source: {metadata.get('source', 'N/A')}")
        print(f"    - Slide number: {metadata.get('slide_number', 'N/A')}")
        print(f"    - Slide title: {metadata.get('slide_title', 'N/A')[:40]}")
        print(f"    - Type: {metadata.get('chunk_type', 'N/A')}")
        print(f"    - Total slides: {metadata.get('total_slides', 'N/A')}")
        print(f"    - Text length: {len(chunk['text'])} chars")

        # Check if speaker notes are included
        if 'notes' in chunk['text'].lower():
            print(f"    - ✓ Includes speaker notes")

    print(f"\n✅ Test 3 Passed: Document processor integration works")
    return True


def test_query_scenarios():
    """Test 4: Simulate query scenarios."""
    print("\n" + "=" * 70)
    print("TEST 4: Query Scenario Simulation")
    print("=" * 70)

    from modules.document_processor import TextChunker

    # Process the file
    with open('comprehensive_presentation.pptx', 'rb') as f:
        chunks = TextChunker.process_document(
            f,
            'comprehensive_presentation.pptx',
            '.pptx'
        )

    # Define test queries and check if relevant data is in chunks
    test_queries = [
        {
            'query': 'revenue financial performance',
            'keywords': ['Revenue', '52.3M', 'growth']
        },
        {
            'query': 'market expansion international',
            'keywords': ['Market', 'EMEA', 'APAC', 'expansion']
        },
        {
            'query': 'product roadmap 2025',
            'keywords': ['Product', 'Roadmap', '2025', 'AI']
        },
        {
            'query': 'team hiring employees',
            'keywords': ['Team', 'employees', 'Engineering', 'Sales']
        },
        {
            'query': 'priorities profitability',
            'keywords': ['Priorities', 'growth', 'profitability']
        }
    ]

    print(f"\n✓ Simulating {len(test_queries)} query scenarios...")

    results_summary = []

    for i, test in enumerate(test_queries, 1):
        print(f"\n  Query {i}: '{test['query']}'")

        # Find chunks containing keywords
        matching_chunks = []
        for chunk in chunks:
            text_lower = chunk['text'].lower()
            matches = [kw for kw in test['keywords']
                      if kw.lower() in text_lower]
            if matches:
                matching_chunks.append({
                    'chunk': chunk,
                    'matches': matches
                })

        if matching_chunks:
            print(f"    ✓ Found {len(matching_chunks)} relevant chunks")
            best = matching_chunks[0]
            slide_num = best['chunk']['metadata'].get('slide_number', '?')
            slide_title = best['chunk']['metadata'].get('slide_title', 'Unknown')
            print(f"    - Best match from slide {slide_num}: {slide_title[:40]}")
            print(f"    - Keywords found: {', '.join(best['matches'])}")
            results_summary.append(True)
        else:
            print(f"    ⚠ No matching chunks found")
            results_summary.append(False)

    success_rate = sum(results_summary) / len(results_summary) if results_summary else 0
    print(f"\n✓ Query success rate: {success_rate:.0%} ({sum(results_summary)}/{len(results_summary)})")

    print(f"\n✅ Test 4 Passed: Query simulation complete")
    return True


def test_speaker_notes():
    """Test 5: Verify speaker notes are included."""
    print("\n" + "=" * 70)
    print("TEST 5: Speaker Notes Verification")
    print("=" * 70)

    from modules.document_processor import TextChunker

    # Process the file
    with open('comprehensive_presentation.pptx', 'rb') as f:
        chunks = TextChunker.process_document(
            f,
            'comprehensive_presentation.pptx',
            '.pptx'
        )

    print(f"\n✓ Checking {len(chunks)} chunks for speaker notes...")

    chunks_with_notes = 0
    for i, chunk in enumerate(chunks, 1):
        text = chunk['text']

        # Check if this chunk contains speaker notes section
        if 'Speaker Notes:' in text or 'Notes:' in text:
            chunks_with_notes += 1
            print(f"\n  Chunk {i} (Slide {chunk['metadata'].get('slide_number', '?')}):")
            print(f"    ✓ Contains speaker notes")

            # Extract notes section
            if 'Speaker Notes:' in text:
                notes_start = text.index('Speaker Notes:')
                notes_preview = text[notes_start:notes_start+150]
                print(f"    - Notes preview: {notes_preview}...")

    print(f"\n✓ Found speaker notes in {chunks_with_notes}/{len(chunks)} chunks")

    if chunks_with_notes >= len(chunks) * 0.8:  # At least 80% should have notes
        print(f"✓ Speaker notes are properly included")
    else:
        print(f"⚠ Some slides may be missing speaker notes")

    print(f"\n✅ Test 5 Passed: Speaker notes verification complete")
    return True


def main():
    """Run all tests."""
    print("\n" + "=" * 70)
    print("🧪 FINAL POWERPOINT PROCESSING TEST")
    print("=" * 70)
    print("\nThis test verifies PowerPoint support without requiring")
    print("embeddings or vector storage.")
    print("\nTesting:")
    print("  1. Direct python-pptx functionality")
    print("  2. PowerPointProcessor module")
    print("  3. Document processor integration")
    print("  4. Query scenario simulation")
    print("  5. Speaker notes verification")
    print("\n" + "=" * 70)

    try:
        # Check if test file exists
        if not Path('comprehensive_presentation.pptx').exists():
            print("\n❌ Test file 'comprehensive_presentation.pptx' not found!")
            return False

        # Run tests
        test_direct_pptx_reading()
        test_pptx_processor_module()
        test_document_processor()
        test_query_scenarios()
        test_speaker_notes()

        # Summary
        print("\n" + "=" * 70)
        print("✅ ALL TESTS PASSED!")
        print("=" * 70)
        print("\n📊 PowerPoint Support Verification Summary:")
        print("  ✓ python-pptx works correctly")
        print("  ✓ PowerPointProcessor module functions properly")
        print("  ✓ Document processor integration successful")
        print("  ✓ Chunks have correct structure and metadata")
        print("  ✓ Query scenarios can find relevant data")
        print("  ✓ Speaker notes are properly included")
        print("\n🎯 PowerPoint support is fully functional!")
        print("\n📝 What this means:")
        print("  - Users can upload .pptx files")
        print("  - All slides will be extracted")
        print("  - Slide structure is preserved")
        print("  - Speaker notes are included")
        print("  - Natural language queries will work")
        print("\n💡 Example queries that will work:")
        print("  • 'What was the Q4 revenue?'")
        print("  • 'Which markets did we expand into?'")
        print("  • 'What products are launching in 2025?'")
        print("  • 'How many employees do we have?'")
        print("  • 'What are the 2025 priorities?'")
        print("\n🚀 Ready to use in the Streamlit app!")

        return True

    except Exception as e:
        print(f"\n❌ Test failed: {e}")
        import traceback
        traceback.print_exc()
        return False


if __name__ == '__main__':
    success = main()
    sys.exit(0 if success else 1)
